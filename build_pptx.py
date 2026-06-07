# -*- coding: utf-8 -*-
"""
计算机视觉发展史 · 简约风格 PPT (Schneider Green 主题)
- 12 slides: Hero + 6 Acts + 数据集 + 应用 + 关键人物 + 收尾
- 16:9 widescreen
- 输出: cv_history_deck.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# =============== 主题色 · 施耐德绿 (Schneider Green) ===============
SCH_GREEN      = RGBColor(0x3D, 0xCD, 0x58)  # 施耐德绿主色
SCH_GREEN_DARK = RGBColor(0x00, 0x75, 0x4A)  # 深绿(暗底使用)
SCH_GREEN_TINT = RGBColor(0xE8, 0xF8, 0xEE)  # 极浅绿底
SCH_GREEN_SOFT = RGBColor(0x2D, 0xAA, 0x48)  # 柔和绿
INK            = RGBColor(0x0A, 0x0A, 0x0A)  # 主文字
PAPER          = RGBColor(0xFA, 0xFA, 0xF8)  # 主底色
GREY_BG        = RGBColor(0xF5, 0xF5, 0xF3)  # 浅灰底
GREY_RULE      = RGBColor(0xD4, 0xD4, 0xD2)  # 分割线
GREY_MID       = RGBColor(0x73, 0x73, 0x73)  # 辅助文字
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)

# =============== 字体 (跨平台安全) ===============
FONT_HEAD = "Microsoft YaHei"   # 中文标题
FONT_BODY = "Microsoft YaHei UI"  # 中文正文
FONT_MONO = "Consolas"          # 等宽

# =============== 尺寸 ===============
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.6)

# =============== 初始化 Presentation ===============
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]
print("Build start")

# =============== 辅助函数 ===============
def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None, shadow=False):
    """添加矩形"""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_w is not None:
            shp.line.width = line_w
    return shp

def add_line(slide, x1, y1, x2, y2, color=None, weight=Pt(0.75)):
    """添加直线 (横线/竖线)"""
    from pptx.util import Emu
    shp = slide.shapes.add_connector(1, x1, y1, x2, y2)
    shp.line.color.rgb = color if color else GREY_RULE
    shp.line.width = weight
    return shp

def add_text(slide, x, y, w, h, text, *, font=FONT_BODY, size=14, bold=False,
             color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
             line_spacing=1.2):
    """添加文本框 (单段多 run 由换行符 \n 切分)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb

def add_multi(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    """runs = list of dict: {text, font, size, bold, italic, color}
    runs 内可通过 'newline': True 触发换段"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    first = True
    for spec in runs:
        if spec.get("newline"):
            p = tf.add_paragraph()
            p.alignment = align
            p.line_spacing = line_spacing
            first = True
            continue
        r = p.add_run()
        r.text = spec["text"]
        r.font.name = spec.get("font", FONT_BODY)
        r.font.size = Pt(spec.get("size", 14))
        r.font.bold = spec.get("bold", False)
        r.font.italic = spec.get("italic", False)
        r.font.color.rgb = spec.get("color", INK)
        first = False
    return tb

def page_chrome(slide, idx, total, kicker="", show_nav=True):
    """统一 chrome: 顶 meta + 底脚 (页码 + 主题名)"""
    # 顶部: 左 kicker + 右 SCHNEIDER GREEN · CV HISTORY
    add_text(slide, MARGIN, Inches(0.35), Inches(6), Inches(0.3),
             kicker, font=FONT_MONO, size=10, color=GREY_MID)
    add_text(slide, SLIDE_W - MARGIN - Inches(5), Inches(0.35), Inches(5), Inches(0.3),
             "SCHNEIDER GREEN  ·  CV HISTORY", font=FONT_MONO, size=10,
             color=GREY_MID, align=PP_ALIGN.RIGHT)
    # 顶部分割线
    add_line(slide, MARGIN, Inches(0.72), SLIDE_W - MARGIN, Inches(0.72),
             color=GREY_RULE, weight=Pt(0.5))
    # 底部: 左签名 + 右页码
    add_line(slide, MARGIN, SLIDE_H - Inches(0.55), SLIDE_W - MARGIN, SLIDE_H - Inches(0.55),
             color=GREY_RULE, weight=Pt(0.5))
    add_text(slide, MARGIN, SLIDE_H - Inches(0.45), Inches(8), Inches(0.3),
             "计算机视觉发展史  ·  1963 → 2024", font=FONT_MONO, size=10, color=GREY_MID)
    if show_nav:
        add_text(slide, SLIDE_W - MARGIN - Inches(3), SLIDE_H - Inches(0.45),
                 Inches(3), Inches(0.3),
                 f"{idx:02d} / {total:02d}", font=FONT_MONO, size=10,
                 color=GREY_MID, align=PP_ALIGN.RIGHT)

def accent_bar(slide, x, y, h, w=Inches(0.08), color=SCH_GREEN):
    """左侧施耐德绿 accent 竖条"""
    add_rect(slide, x, y, w, h, fill=color)

print("Helpers loaded")

# =============== Slide 1 · HERO 封面 ===============
def slide_1_hero():
    s = prs.slides.add_slide(BLANK)
    # 全幅深绿底
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=SCH_GREEN_DARK)
    # 右上: 60 年数字大标
    add_text(s, SLIDE_W - Inches(5.5), Inches(0.4), Inches(5), Inches(2.6),
             "60", font=FONT_HEAD, size=160, bold=False, color=SCH_GREEN,
             align=PP_ALIGN.RIGHT)
    add_text(s, SLIDE_W - Inches(5.5), Inches(2.4), Inches(5), Inches(0.4),
             "YEARS  ·  1963 — 2024", font=FONT_MONO, size=14,
             color=SCH_GREEN, align=PP_ALIGN.RIGHT)
    # 大标题 (左下)
    add_text(s, MARGIN, Inches(2.6), Inches(9), Inches(0.5),
             "看见的六十年", font=FONT_MONO, size=18, color=WHITE)
    add_multi(s, MARGIN, Inches(3.1), Inches(11), Inches(2.6), [
        {"text": "计算机视觉", "font": FONT_HEAD, "size": 96, "bold": False, "color": WHITE},
        {"newline": True},
        {"text": "发展史", "font": FONT_HEAD, "size": 96, "bold": False, "color": SCH_GREEN},
    ], line_spacing=1.0)
    # 副标
    add_text(s, MARGIN, Inches(5.5), Inches(11), Inches(0.5),
             "From Pixels to Perception · 从像素到认知",
             font=FONT_BODY, size=20, color=WHITE, italic=True)
    # 装饰: 底部细线
    add_line(s, MARGIN, SLIDE_H - Inches(1.0), SLIDE_W - MARGIN, SLIDE_H - Inches(1.0),
             color=SCH_GREEN, weight=Pt(1.5))
    add_text(s, MARGIN, SLIDE_H - Inches(0.7), Inches(6), Inches(0.3),
             "线下分享  ·  2026.06.07", font=FONT_MONO, size=12, color=WHITE)
    add_text(s, SLIDE_W - MARGIN - Inches(4), SLIDE_H - Inches(0.7), Inches(4), Inches(0.3),
             "HERO  ·  01 / 14", font=FONT_MONO, size=12, color=SCH_GREEN,
             align=PP_ALIGN.RIGHT)

print("Slide 1 defined")

# =============== Slide 2 · 总览 · 6 个时代 ===============
def slide_2_overview():
    s = prs.slides.add_slide(BLANK)
    # 白底
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=PAPER)
    page_chrome(s, 2, 14, kicker="OVERVIEW  ·  6 ACTS  ·  60 YEARS")
    # 大标题
    add_text(s, MARGIN, Inches(1.1), Inches(8), Inches(0.5),
             "THE  FIELD  IN  SIX  ACTS",
             font=FONT_MONO, size=14, color=SCH_GREEN_SOFT, bold=True)
    add_text(s, MARGIN, Inches(1.5), Inches(11), Inches(1.3),
             "六幕  ·  视觉的演化",
             font=FONT_HEAD, size=54, bold=False, color=INK)
    add_text(s, MARGIN, Inches(2.7), Inches(11), Inches(0.5),
             "从符号推理到多模态大模型,一条跨越六十年的认知曲线。",
             font=FONT_BODY, size=18, color=GREY_MID)
    # 6 个时代卡片
    acts = [
        ("01", "1960s–1970s", "奠基期", "积木世界 · Marr 三层"),
        ("02", "1980s–1990s", "几何视觉", "SIFT · HOG · 多视几何"),
        ("03", "2000s", "统计学习", "ImageNet · BoVW · DPM"),
        ("04", "2012–2017", "深度学习", "AlexNet · ResNet · YOLO"),
        ("05", "2018–2020", "生成与注意力", "GAN · ViT · NeRF"),
        ("06", "2021 →", "多模态", "CLIP · SAM · Sora"),
    ]
    card_w = (SLIDE_W - MARGIN * 2 - Inches(0.4) * 5) / 6
    card_h = Inches(3.0)
    card_y = Inches(3.6)
    for i, (nb, era, title, sub) in enumerate(acts):
        x = MARGIN + i * (card_w + Inches(0.4))
        # 卡片白底
        add_rect(s, x, card_y, card_w, card_h, fill=WHITE, line=GREY_RULE, line_w=Pt(0.5))
        # 顶部 8px 施耐德绿条
        add_rect(s, x, card_y, card_w, Inches(0.08), fill=SCH_GREEN)
        # 编号
        add_text(s, x + Inches(0.25), card_y + Inches(0.3), card_w - Inches(0.5), Inches(0.6),
                 nb, font=FONT_HEAD, size=36, color=SCH_GREEN, bold=False)
        # 年代
        add_text(s, x + Inches(0.25), card_y + Inches(1.0), card_w - Inches(0.5), Inches(0.3),
                 era, font=FONT_MONO, size=11, color=GREY_MID)
        # 标题
        add_text(s, x + Inches(0.25), card_y + Inches(1.4), card_w - Inches(0.5), Inches(0.6),
                 title, font=FONT_HEAD, size=20, color=INK, bold=True)
        # 副标
        add_text(s, x + Inches(0.25), card_y + Inches(2.1), card_w - Inches(0.5), Inches(0.8),
                 sub, font=FONT_BODY, size=11, color=GREY_MID, line_spacing=1.4)

print("Slide 2 defined")

# =============== 通用 · 时代章节页(act divider) ===============
def act_divider(act_nb, era, title, subtitle, keypoints, slide_idx, total=14, dark=False):
    """Act 章节封面页: 大数字 + 时代 + 主题 + 关键节点"""
    s = prs.slides.add_slide(BLANK)
    bg = INK if dark else PAPER
    fg = PAPER if dark else INK
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=bg)
    # chrome
    add_text(s, MARGIN, Inches(0.35), Inches(6), Inches(0.3),
             f"ACT  {act_nb}  ·  {era}", font=FONT_MONO, size=10,
             color=SCH_GREEN if dark else SCH_GREEN_SOFT, bold=True)
    add_text(s, SLIDE_W - MARGIN - Inches(5), Inches(0.35), Inches(5), Inches(0.3),
             f"CHAPTER  ·  {slide_idx:02d} / {total:02d}", font=FONT_MONO, size=10,
             color=GREY_MID, align=PP_ALIGN.RIGHT)
    add_line(s, MARGIN, Inches(0.72), SLIDE_W - MARGIN, Inches(0.72),
             color=SCH_GREEN, weight=Pt(1.0))
    # 巨大 act 编号
    add_text(s, MARGIN, Inches(1.0), Inches(7), Inches(3.0),
             act_nb, font=FONT_HEAD, size=180, color=SCH_GREEN, bold=False)
    # 时代
    add_text(s, MARGIN, Inches(4.4), Inches(11), Inches(0.5),
             era, font=FONT_MONO, size=14, color=fg)
    # 大标题
    add_text(s, MARGIN, Inches(4.8), Inches(11), Inches(1.0),
             title, font=FONT_HEAD, size=46, color=fg, bold=True)
    # 副标
    add_text(s, MARGIN, Inches(5.6), Inches(11), Inches(0.6),
             subtitle, font=FONT_BODY, size=18, color=fg, italic=True, line_spacing=1.4)
    # 关键节点 (右下)
    if keypoints:
        kx = SLIDE_W - Inches(4.6)
        ky = Inches(1.4)
        add_text(s, kx, ky, Inches(4), Inches(0.3),
                 "KEY POINTS", font=FONT_MONO, size=10, color=SCH_GREEN, bold=True)
        add_line(s, kx, ky + Inches(0.35), kx + Inches(4), ky + Inches(0.35),
                 color=SCH_GREEN, weight=Pt(0.75))
        for i, kp in enumerate(keypoints):
            add_text(s, kx, ky + Inches(0.55) + Inches(0.5) * i, Inches(0.4), Inches(0.4),
                     f"{i+1:02d}", font=FONT_MONO, size=11, color=SCH_GREEN, bold=True)
            add_text(s, kx + Inches(0.5), ky + Inches(0.55) + Inches(0.5) * i, Inches(3.5), Inches(0.4),
                     kp, font=FONT_BODY, size=13, color=fg)
    return s

def act_content(act_nb, era, theme, slide_idx, items, total=14, dark=False):
    """Act 内容页: 顶部 chrome + 大主题 + 网格 items"""
    s = prs.slides.add_slide(BLANK)
    bg = GREY_BG if not dark else INK
    fg = INK if not dark else PAPER
    sub_fg = GREY_MID if not dark else RGBColor(0xC0, 0xC0, 0xC0)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=bg)
    # chrome
    add_text(s, MARGIN, Inches(0.35), Inches(6), Inches(0.3),
             f"ACT  {act_nb}  ·  {era}", font=FONT_MONO, size=10,
             color=SCH_GREEN_SOFT, bold=True)
    add_text(s, SLIDE_W - MARGIN - Inches(5), Inches(0.35), Inches(5), Inches(0.3),
             f"CONTENT  ·  {slide_idx:02d} / {total:02d}", font=FONT_MONO, size=10,
             color=sub_fg, align=PP_ALIGN.RIGHT)
    add_line(s, MARGIN, Inches(0.72), SLIDE_W - MARGIN, Inches(0.72),
             color=SCH_GREEN, weight=Pt(0.75))
    # 大主题
    add_text(s, MARGIN, Inches(0.95), Inches(11), Inches(0.4),
             theme, font=FONT_HEAD, size=32, color=fg, bold=True)
    # 主题下绿条
    add_rect(s, MARGIN, Inches(1.55), Inches(0.6), Inches(0.06), fill=SCH_GREEN)
    # items 网格 (3列)
    col_w = (SLIDE_W - MARGIN * 2 - Inches(0.4) * 2) / 3
    row_h = Inches(2.3)
    grid_y = Inches(1.95)
    for i, item in enumerate(items):
        col = i % 3
        row = i // 3
        x = MARGIN + col * (col_w + Inches(0.4))
        y = grid_y + row * (row_h + Inches(0.2))
        # 卡片
        add_rect(s, x, y, col_w, row_h, fill=PAPER if not dark else RGBColor(0x1A, 0x1A, 0x1A),
                 line=GREY_RULE if not dark else RGBColor(0x33, 0x33, 0x33), line_w=Pt(0.5))
        # 左侧 accent 竖条
        add_rect(s, x, y, Inches(0.08), row_h, fill=SCH_GREEN)
        # 年份
        add_text(s, x + Inches(0.3), y + Inches(0.2), col_w - Inches(0.5), Inches(0.3),
                 item["year"], font=FONT_MONO, size=11, color=SCH_GREEN, bold=True)
        # 名字
        add_text(s, x + Inches(0.3), y + Inches(0.5), col_w - Inches(0.5), Inches(0.5),
                 item["name"], font=FONT_HEAD, size=20,
                 color=fg if not dark else PAPER, bold=True)
        # 描述
        add_text(s, x + Inches(0.3), y + Inches(1.05), col_w - Inches(0.5), row_h - Inches(1.2),
                 item["desc"], font=FONT_BODY, size=12,
                 color=sub_fg, line_spacing=1.4)
    return s

print("Act helpers loaded")

# =============== Slide 3 · Act I 章节封面 (奠基期) ===============
def slide_3_act1():
    act_divider(
        "I", "1960s — 1970s", "奠基期",
        "学科诞生 · 积木世界 · Marr 三层框架",
        ["1963  Larry Roberts 博士论文", "1966  MIT Summer Vision Project",
         "1970s  Marr 视觉计算理论", "1986  Canny 最优边缘检测"],
        slide_idx=3
    )

# =============== Slide 4 · Act I 内容页 ===============
def slide_4_act1_content():
    items = [
        {"year": "1963", "name": "Larry Roberts", "desc": "MIT 博士论文《Machine Perception of Three-Dimensional Solids》,从 2D 线条图推断 3D 几何——计算机视觉的奠基之作。"},
        {"year": "1966", "name": "Summer Vision", "desc": "MIT Papert 教授启动「Summer Vision Project」,尝试一个夏天解决视觉问题。目标远未达成,但确立了「视觉可被计算建模」的范式。"},
        {"year": "1970s", "name": "David Marr", "desc": "MIT AI Lab 提出视觉计算三层框架:计算理论层(What & Why) / 算法表示层(How) / 物理实现层(硬件)。"},
        {"year": "1965/68", "name": "Roberts · Sobel", "desc": "最早的边缘检测算子。Roberts 算子(1965)、Sobel 算子(1968)开启视觉特征提取时代。"},
        {"year": "1970", "name": "Prewitt 算子", "desc": "基于梯度近似的边缘检测,与 Sobel 同属一脉,至今仍是图像处理教材的入门内容。"},
        {"year": "1981", "name": "Horn · Lucas-Kanade", "desc": "光流估计的两篇里程碑论文,奠定了运动视觉的数学基础,后续三十年的跟踪算法均从中衍生。"},
    ]
    act_content("I", "1960s — 1970s  ·  奠基期", "当机器开始「看见」", 4, items)

# =============== Slide 5 · Act II 章节封面 (几何视觉) ===============
def slide_5_act2():
    act_divider(
        "II", "1980s — 1990s", "几何视觉时代",
        "多视几何 · 特征工程 · 早期检测",
        ["1988  Harris 角点检测", "1999  SIFT 尺度不变特征",
         "2001  Viola-Jones 人脸检测", "2005  HOG + SVM 行人检测"],
        slide_idx=5, dark=True
    )

# =============== Slide 6 · Act II 内容页 ===============
def slide_6_act2_content():
    items = [
        {"year": "1981 / 1992", "name": "多视几何", "desc": "Longuet-Higgins 8 点法(1981)与 Hartley & Zisserman《Multiple View Geometry》(1992),确立从 2D 图像到 3D 重建的数学基础。"},
        {"year": "1987 / 2000", "name": "相机标定", "desc": "Tsai(1987)与张正友(2000)标定法,让相机内外参可被精确估计,3D 视觉走向工程化。"},
        {"year": "1988", "name": "Harris 角点", "desc": "Chris Harris & Mike Stephens 提出第一个广泛使用的特征检测器,基于二阶矩的角点响应函数。"},
        {"year": "1999", "name": "SIFT", "desc": "David Lowe 提出尺度不变特征变换(SIFT),对旋转、尺度、光照保持鲁棒,是特征点时代的标杆,被引用超 7 万次。"},
        {"year": "2001", "name": "Viola-Jones", "desc": "Haar 特征 + AdaBoost + 级联分类器,首次让实时人脸检测跑在 PC 上,宣告「检测」成为独立任务。"},
        {"year": "2005", "name": "HOG + SVM", "desc": "Dalal & Triggs 的方向梯度直方图,行人检测的经典组合,启发了后续十多年的传统检测器。"},
    ]
    act_content("II", "1980s — 1990s  ·  几何视觉", "用手工特征丈量世界", 6, items)

print("Act I + II defined")

# =============== Slide 7 · Act III 章节封面 (统计学习) ===============
def slide_7_act3():
    act_divider(
        "III", "2000s", "统计学习崛起",
        "数据集革命 · 视觉词袋 · DPM",
        ["1998  MNIST 手写数字", "2005  PASCAL VOC 基准",
         "2009  ImageNet 14M 图像", "2010  DPM 可变形部件"],
        slide_idx=7
    )

# =============== Slide 8 · Act III 内容页 ===============
def slide_8_act3_content():
    items = [
        {"year": "1998", "name": "MNIST", "desc": "LeCun 等人发布的手写数字数据集,7 万张 28×28 灰度图,被戏称为「机器学习的果蝇」——三十年后仍出现在每个深度学习教程第一章。"},
        {"year": "2003 / 2006", "name": "Caltech 101 / 256", "desc": "早期图像分类基准,9 千张图像、102 类,让「分类」成为 CV 任务的度量衡。"},
        {"year": "2005 — 2012", "name": "PASCAL VOC", "desc": "第一个真正意义上的多类别检测 / 分割基准,20 类,催生了大量经典检测算法。"},
        {"year": "2009", "name": "ImageNet", "desc": "Fei-Fei Li 主导,1400 万张标注图像、2.2 万类,数据集规模直接拉升三个数量级——深度学习革命的真正火种。"},
        {"year": "2006", "name": "空间金字塔匹配", "desc": "Lazebnik 等把空间金字塔引入 BoVW,多尺度特征融合提升分类精度。"},
        {"year": "2010", "name": "DPM", "desc": "Felzenszwalb 等提出的可变形部件模型,连续多年 PASCAL VOC 冠军,是深度学习前传统检测的巅峰。"},
    ]
    act_content("III", "2000s  ·  统计学习", "数据为王  ·  特征为器", 8, items)

# =============== Slide 9 · Act IV 章节封面 (深度学习) ===============
def slide_9_act4():
    act_divider(
        "IV", "2012 — 2017", "深度学习革命",
        "AlexNet 时刻 · CNN 架构 · 检测分割",
        ["2012  AlexNet 15.3% Top-5 错误率", "2015  ResNet 152 层",
         "2016  YOLO 实时检测", "2017  Mask R-CNN 实例分割"],
        slide_idx=9, dark=True
    )

# =============== Slide 10 · Act IV 内容页 ===============
def slide_10_act4_content():
    items = [
        {"year": "2012", "name": "AlexNet", "desc": "Krizhevsky / Sutskever / Hinton 以 15.3% Top-5 错误率夺冠 ImageNet(第二名 26.2%),ReLU + Dropout + GPU 训练三件套,深度卷积正式登场。"},
        {"year": "2014", "name": "VGG · GoogLeNet", "desc": "VGG 用 3×3 小卷积堆到 19 层,GoogLeNet 提出 Inception 多尺度模块。两条路径都指向「更深、更宽」。"},
        {"year": "2015", "name": "ResNet", "desc": "何恺明等提出残差连接,网络深度跃至 152 层,解决退化问题,横扫 ImageNet 与 COCO。"},
        {"year": "2015", "name": "Faster R-CNN · FCN", "desc": "Faster R-CNN 用 RPN 统一两阶段检测;Long 等提出 FCN,端到端像素分类开启语义分割时代。"},
        {"year": "2016", "name": "YOLO · SSD", "desc": "Redmon 等的 YOLO 端到端单阶段检测,达到实时;SSD 用多尺度特征图。检测从此「快」与「准」可兼得。"},
        {"year": "2017", "name": "Mask R-CNN · DenseNet", "desc": "何恺明等在 Faster R-CNN 基础上加掩码分支,实例分割的标杆;同年 DenseNet 探索极致密集连接。"},
    ]
    act_content("IV", "2012 — 2017  ·  深度学习", "像素开始可被学习", 10, items)

print("Act III + IV defined")

# =============== Slide 11 · Act V 章节封面 + 内容合并页 (生成/Transformer) ===============
def slide_11_act5():
    """Act V 是范式拐点,做成 1.5 页: 左侧大引述 + 右侧三个时间线里程碑"""
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=INK)
    # chrome
    add_text(s, MARGIN, Inches(0.35), Inches(6), Inches(0.3),
             "ACT  V  ·  2018 — 2020", font=FONT_MONO, size=10,
             color=SCH_GREEN, bold=True)
    add_text(s, SLIDE_W - MARGIN - Inches(5), Inches(0.35), Inches(5), Inches(0.3),
             "PARADIGM  SHIFT  ·  11 / 14", font=FONT_MONO, size=10,
             color=RGBColor(0xC0, 0xC0, 0xC0), align=PP_ALIGN.RIGHT)
    add_line(s, MARGIN, Inches(0.72), SLIDE_W - MARGIN, Inches(0.72),
             color=SCH_GREEN, weight=Pt(1.0))
    # 左侧: 大引述
    add_text(s, MARGIN, Inches(1.2), Inches(7.5), Inches(0.5),
             "PARADIGM  SHIFT", font=FONT_MONO, size=14, color=SCH_GREEN, bold=True)
    add_multi(s, MARGIN, Inches(1.6), Inches(7.5), Inches(4.5), [
        {"text": "从", "font": FONT_HEAD, "size": 52, "color": WHITE},
        {"text": " 像素", "font": FONT_HEAD, "size": 52, "color": WHITE, "italic": True},
        {"text": " 到 ", "font": FONT_HEAD, "size": 52, "color": WHITE},
        {"text": "token", "font": FONT_HEAD, "size": 52, "color": SCH_GREEN},
        {"text": "。", "font": FONT_HEAD, "size": 52, "color": WHITE},
        {"newline": True},
        {"text": "从", "font": FONT_HEAD, "size": 52, "color": WHITE},
        {"text": " 2D", "font": FONT_HEAD, "size": 52, "color": WHITE, "italic": True},
        {"text": " 到 ", "font": FONT_HEAD, "size": 52, "color": WHITE},
        {"text": "3D", "font": FONT_HEAD, "size": 52, "color": SCH_GREEN},
        {"text": "。", "font": FONT_HEAD, "size": 52, "color": WHITE},
        {"newline": True},
        {"text": "从", "font": FONT_HEAD, "size": 52, "color": WHITE},
        {"text": " 看见", "font": FONT_HEAD, "size": 52, "color": WHITE, "italic": True},
        {"text": " 到 ", "font": FONT_HEAD, "size": 52, "color": WHITE},
        {"text": "生成", "font": FONT_HEAD, "size": 52, "color": SCH_GREEN},
        {"text": "。", "font": FONT_HEAD, "size": 52, "color": WHITE},
    ], line_spacing=1.15)
    # 右侧: 三个里程碑
    milestones = [
        ("2014", "GAN", "Goodfellow 提出生成对抗网络,从此「生成」成为独立任务。"),
        ("2020", "ViT", "Dosovitskiy 等用纯 Transformer 在 ImageNet 上超越 CNN。"),
        ("2020", "NeRF", "Mildenhall 等用 MLP 隐式表示 3D 场景,3D 视觉新范式。"),
    ]
    rx = SLIDE_W - Inches(4.6)
    ry = Inches(1.2)
    add_text(s, rx, ry, Inches(4), Inches(0.3),
             "MILESTONES", font=FONT_MONO, size=10, color=SCH_GREEN, bold=True)
    add_line(s, rx, ry + Inches(0.35), rx + Inches(4), ry + Inches(0.35),
             color=SCH_GREEN, weight=Pt(0.75))
    for i, (yr, name, desc) in enumerate(milestones):
        y0 = ry + Inches(0.55) + Inches(1.5) * i
        add_text(s, rx, y0, Inches(1.0), Inches(0.4),
                 yr, font=FONT_MONO, size=12, color=SCH_GREEN, bold=True)
        add_text(s, rx + Inches(1.0), y0, Inches(3), Inches(0.4),
                 name, font=FONT_HEAD, size=24, color=WHITE, bold=True)
        add_text(s, rx, y0 + Inches(0.55), Inches(4), Inches(0.9),
                 desc, font=FONT_BODY, size=12, color=RGBColor(0xC0, 0xC0, 0xC0), line_spacing=1.4)
    # 底部 takeaway
    add_line(s, MARGIN, SLIDE_H - Inches(1.0), SLIDE_W - MARGIN, SLIDE_H - Inches(1.0),
             color=WHITE, weight=Pt(0.5))
    add_text(s, MARGIN, SLIDE_H - Inches(0.8), Inches(11), Inches(0.5),
             "三件事:  Attention 替代了卷积  ·  隐式表示重写了 3D  ·  生成开始反向理解感知。",
             font=FONT_BODY, size=14, color=WHITE, italic=True)

# =============== Slide 12 · Act VI 内容页 (多模态与基础模型) ===============
def slide_12_act6():
    """Act VI 内容页: 多模态与基础模型时代"""
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=GREY_BG)
    # chrome
    add_text(s, MARGIN, Inches(0.35), Inches(6), Inches(0.3),
             "ACT  VI  ·  2021 → NOW", font=FONT_MONO, size=10,
             color=SCH_GREEN_SOFT, bold=True)
    add_text(s, SLIDE_W - MARGIN - Inches(5), Inches(0.35), Inches(5), Inches(0.3),
             "FOUNDATION  MODELS  ·  12 / 14", font=FONT_MONO, size=10,
             color=GREY_MID, align=PP_ALIGN.RIGHT)
    add_line(s, MARGIN, Inches(0.72), SLIDE_W - MARGIN, Inches(0.72),
             color=SCH_GREEN, weight=Pt(0.75))
    # 大主题
    add_text(s, MARGIN, Inches(0.95), Inches(11), Inches(0.4),
             "多模态与基础模型", font=FONT_HEAD, size=32, color=INK, bold=True)
    add_text(s, MARGIN, Inches(1.5), Inches(11), Inches(0.5),
             "When vision learned to read, write, and dream.",
             font=FONT_BODY, size=16, color=GREY_MID, italic=True)
    add_rect(s, MARGIN, Inches(2.05), Inches(0.6), Inches(0.06), fill=SCH_GREEN)
    # items 网格 2x3
    items = [
        {"year": "2021", "name": "CLIP", "desc": "OpenAI 用 4 亿图文对训练,零样本分类能力惊人,首次让视觉与语言进入同一表征空间。"},
        {"year": "2022", "name": "Stable Diffusion", "desc": "潜在扩散模型 + 开放权重,催生 AIGC 浪潮。「画」从专家技能变成一句话的事。"},
        {"year": "2022 / 23", "name": "DALL·E 2 · Midjourney", "desc": "文本生成图像进入消费级市场,Imagen 与 Midjourney v5 把真实感推到新高度。"},
        {"year": "2023", "name": "SAM · GPT-4V", "desc": "Meta 的 Segment Anything 11 亿 mask 通用分割;GPT-4V 让多模态对话成为日常。"},
        {"year": "2023", "name": "3D Gaussian Splatting", "desc": "Kerbl 等用高斯泼溅做实时高质量新视角合成,SIGGRAPH 最佳论文,NeRF 之后新一代 3D 表征。"},
        {"year": "2024", "name": "Sora · 具身智能", "desc": "OpenAI 的 Sora 用扩散 Transformer 生成视频;RT-2 / π0 把视觉-语言-动作闭环推向真实物理世界。"},
    ]
    col_w = (SLIDE_W - MARGIN * 2 - Inches(0.4) * 2) / 3
    row_h = Inches(1.95)
    grid_y = Inches(2.4)
    for i, item in enumerate(items):
        col = i % 3
        row = i // 3
        x = MARGIN + col * (col_w + Inches(0.4))
        y = grid_y + row * (row_h + Inches(0.2))
        add_rect(s, x, y, col_w, row_h, fill=PAPER, line=GREY_RULE, line_w=Pt(0.5))
        add_rect(s, x, y, Inches(0.08), row_h, fill=SCH_GREEN)
        add_text(s, x + Inches(0.3), y + Inches(0.2), col_w - Inches(0.5), Inches(0.3),
                 item["year"], font=FONT_MONO, size=11, color=SCH_GREEN, bold=True)
        add_text(s, x + Inches(0.3), y + Inches(0.5), col_w - Inches(0.5), Inches(0.5),
                 item["name"], font=FONT_HEAD, size=20, color=INK, bold=True)
        add_text(s, x + Inches(0.3), y + Inches(1.05), col_w - Inches(0.5), row_h - Inches(1.2),
                 item["desc"], font=FONT_BODY, size=12, color=GREY_MID, line_spacing=1.4)

print("Act V + VI defined")

# =============== Slide 13 · 数据集 & 应用 ===============
def slide_13_data_apps(idx=13, total=14):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=PAPER)
    page_chrome(s, idx, total, kicker="DATA  +  APPLICATION  ·  1998 → 2024")
    # 大标题
    add_text(s, MARGIN, Inches(1.0), Inches(8), Inches(0.4),
             "DATA  GROWTH  +  APPLICATION  MAP", font=FONT_MONO, size=12,
             color=SCH_GREEN_SOFT, bold=True)
    add_text(s, MARGIN, Inches(1.4), Inches(11), Inches(0.8),
             "数据集曲线  +  产业版图", font=FONT_HEAD, size=36, color=INK, bold=True)
    # 左半: 数据集曲线 (柱状)
    lx = MARGIN
    ly = Inches(2.7)
    lw = Inches(5.6)
    lh = Inches(3.6)
    add_rect(s, lx, ly, lw, lh, fill=WHITE, line=GREY_RULE, line_w=Pt(0.5))
    add_rect(s, lx, ly, Inches(0.08), lh, fill=SCH_GREEN)
    add_text(s, lx + Inches(0.3), ly + Inches(0.2), lw - Inches(0.5), Inches(0.4),
             "DATASET  GROWTH", font=FONT_MONO, size=10, color=SCH_GREEN, bold=True)
    add_text(s, lx + Inches(0.3), ly + Inches(0.55), lw - Inches(0.5), Inches(0.4),
             "70K → 1.1B  masks", font=FONT_HEAD, size=24, color=INK, bold=True)
    # 柱状图 (4 个柱子, 从矮到高)
    data = [
        ("MNIST",  "1998",  0.06, "70K"),
        ("ImageNet", "2009", 0.20, "14M"),
        ("COCO",   "2014", 0.14, "330K"),
        ("SAM-1B", "2023", 0.50, "1.1B"),
    ]
    bar_area_y = ly + Inches(1.2)
    bar_area_h = Inches(2.0)
    bar_w = (lw - Inches(0.6) - Inches(0.3) * 3) / 4
    for i, (name, yr, h_pct, val) in enumerate(data):
        bx = lx + Inches(0.3) + i * (bar_w + Inches(0.3))
        bh = bar_area_h * h_pct * 2  # 放大一些
        # 柱子
        bar_color = SCH_GREEN if i == 3 else GREY_RULE
        add_rect(s, bx, bar_area_y + bar_area_h - bh, bar_w, bh, fill=bar_color)
        # 值标签
        add_text(s, bx, bar_area_y + bar_area_h - bh - Inches(0.35), bar_w, Inches(0.3),
                 val, font=FONT_HEAD, size=14,
                 color=INK if i != 3 else SCH_GREEN, bold=True, align=PP_ALIGN.CENTER)
        # 名称 + 年
        add_text(s, bx, bar_area_y + bar_area_h + Inches(0.1), bar_w, Inches(0.3),
                 name, font=FONT_HEAD, size=11, color=INK, align=PP_ALIGN.CENTER)
        add_text(s, bx, bar_area_y + bar_area_h + Inches(0.4), bar_w, Inches(0.3),
                 yr, font=FONT_MONO, size=10, color=GREY_MID, align=PP_ALIGN.CENTER)
    # 右半: 应用场景列表
    rx = MARGIN + lw + Inches(0.4)
    rw = SLIDE_W - rx - MARGIN
    rh = lh
    add_rect(s, rx, ly, rw, rh, fill=GREY_BG, line=GREY_RULE, line_w=Pt(0.5))
    add_rect(s, rx, ly, Inches(0.08), rh, fill=SCH_GREEN)
    add_text(s, rx + Inches(0.3), ly + Inches(0.2), rw - Inches(0.5), Inches(0.4),
             "APPLICATION  DOMAINS", font=FONT_MONO, size=10, color=SCH_GREEN, bold=True)
    add_text(s, rx + Inches(0.3), ly + Inches(0.55), rw - Inches(0.5), Inches(0.4),
             "八个产业方向", font=FONT_HEAD, size=24, color=INK, bold=True)
    apps = [
        ("人脸识别", "Face++ · 商汤 · 虹软"),
        ("自动驾驶", "Waymo · Tesla · 华为 ADS"),
        ("医疗影像", "肺结节 · 病理切片 · 眼底"),
        ("工业质检", "AOI 视觉 · 表面缺陷"),
        ("AIGC", "SD · Midjourney · Sora"),
        ("AR / VR", "SLAM · 平面检测 · 试穿"),
        ("机器人视觉", "6D Pose · 具身导航"),
        ("遥感", "Land Use · 变化检测"),
    ]
    for i, (name, sub) in enumerate(apps):
        col = i % 2
        row = i // 2
        ax = rx + Inches(0.3) + col * (rw / 2 - Inches(0.15))
        ay = ly + Inches(1.2) + row * Inches(0.5)
        # 圆点
        add_rect(s, ax, ay + Inches(0.1), Inches(0.12), Inches(0.12), fill=SCH_GREEN)
        add_text(s, ax + Inches(0.25), ay, rw / 2 - Inches(0.4), Inches(0.3),
                 name, font=FONT_HEAD, size=13, color=INK, bold=True)
        add_text(s, ax + Inches(0.25), ay + Inches(0.25), rw / 2 - Inches(0.4), Inches(0.3),
                 sub, font=FONT_BODY, size=10, color=GREY_MID)

# =============== Slide 14 · 收尾 / Closing ===============
def slide_14_closing(idx=14, total=14):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=SCH_GREEN_DARK)
    # chrome
    add_text(s, MARGIN, Inches(0.35), Inches(6), Inches(0.3),
             "CLOSING  ·  看见即理解", font=FONT_MONO, size=10,
             color=SCH_GREEN, bold=True)
    add_text(s, SLIDE_W - MARGIN - Inches(5), Inches(0.35), Inches(5), Inches(0.3),
             f"END  ·  {idx:02d} / {total:02d}", font=FONT_MONO, size=10,
             color=WHITE, align=PP_ALIGN.RIGHT)
    add_line(s, MARGIN, Inches(0.72), SLIDE_W - MARGIN, Inches(0.72),
             color=SCH_GREEN, weight=Pt(1.0))
    # 顶部 60
    add_text(s, SLIDE_W - Inches(3.5), Inches(0.8), Inches(3), Inches(2.4),
             "60", font=FONT_HEAD, size=140, color=SCH_GREEN, bold=False,
             align=PP_ALIGN.RIGHT)
    add_text(s, SLIDE_W - Inches(3.5), Inches(2.7), Inches(3), Inches(0.4),
             "YEARS  ·  1963 — 2024", font=FONT_MONO, size=12,
             color=SCH_GREEN, align=PP_ALIGN.RIGHT)
    # 大标题
    add_text(s, MARGIN, Inches(2.0), Inches(9), Inches(0.5),
             "TAKEAWAYS", font=FONT_MONO, size=14, color=SCH_GREEN, bold=True)
    add_multi(s, MARGIN, Inches(2.5), Inches(9.5), Inches(3.5), [
        {"text": "看见的", "font": FONT_HEAD, "size": 80, "color": WHITE},
        {"text": " 六十年", "font": FONT_HEAD, "size": 80, "color": SCH_GREEN, "italic": True},
    ], line_spacing=1.0)
    # 三条 takeaway
    rules = [
        ("01", "从几何到神经",
         "六十年视觉走过了「符号—特征—深度—基础模型」四个范式,每一个都没有否定前一个,而是叠加上去。"),
        ("02", "数据即燃料",
         "从 MNIST 70K 到 SAM-1B 11 亿,数据集曲线是视觉研究最忠实的指标。先有燃料,再有权重。"),
        ("03", "看见即理解",
         "视觉的最终目标不是分类或检测,而是让机器「理解」像素背后的世界。这条路还在继续。"),
    ]
    base_y = Inches(4.0)
    for i, (nb, title, desc) in enumerate(rules):
        x = MARGIN + i * Inches(4.0)
        add_text(s, x, base_y, Inches(3.6), Inches(0.7),
                 nb, font=FONT_HEAD, size=36,
                 color=SCH_GREEN if i == 2 else WHITE, bold=False)
        add_text(s, x, base_y + Inches(0.7), Inches(3.6), Inches(0.4),
                 title, font=FONT_HEAD, size=18, color=WHITE, bold=True)
        add_text(s, x, base_y + Inches(1.1), Inches(3.6), Inches(1.8),
                 desc, font=FONT_BODY, size=12,
                 color=RGBColor(0xE0, 0xE0, 0xE0), line_spacing=1.5)
        # 顶部分割线
        add_line(s, x, base_y - Inches(0.2), x + Inches(3.6), base_y - Inches(0.2),
                 color=WHITE, weight=Pt(0.5))
    # 底部 manifesto
    add_line(s, MARGIN, SLIDE_H - Inches(0.9), SLIDE_W - MARGIN, SLIDE_H - Inches(0.9),
             color=SCH_GREEN, weight=Pt(0.75))
    add_text(s, MARGIN, SLIDE_H - Inches(0.7), Inches(11), Inches(0.4),
             "计算机视觉不只是一门学科,更是连接数字世界与物理世界的桥梁。",
             font=FONT_BODY, size=14, color=WHITE, italic=True, align=PP_ALIGN.CENTER)
    add_text(s, MARGIN, SLIDE_H - Inches(0.35), Inches(11), Inches(0.3),
             "END  OF  FIELD  NOTE  ·  2026.06.07", font=FONT_MONO, size=10,
             color=SCH_GREEN, align=PP_ALIGN.CENTER)

print("Slide 13 + 14 defined")

# =============== 主流程 ===============
TOTAL = 14

# 1. Hero
slide_1_hero()
print("  - Slide 1 (Hero)")

# 2. Overview
slide_2_overview()
print("  - Slide 2 (Overview)")

# 3-4. Act I
slide_3_act1()
print("  - Slide 3 (Act I divider)")
slide_4_act1_content()
print("  - Slide 4 (Act I content)")

# 5-6. Act II
slide_5_act2()
print("  - Slide 5 (Act II divider)")
slide_6_act2_content()
print("  - Slide 6 (Act II content)")

# 7-8. Act III
slide_7_act3()
print("  - Slide 7 (Act III divider)")
slide_8_act3_content()
print("  - Slide 8 (Act III content)")

# 9-10. Act IV
slide_9_act4()
print("  - Slide 9 (Act IV divider)")
slide_10_act4_content()
print("  - Slide 10 (Act IV content)")

# 11. Act V combined
slide_11_act5()
print("  - Slide 11 (Act V)")

# 12. Act VI content
slide_12_act6()
print("  - Slide 12 (Act VI)")

# 13. Data + Apps
slide_13_data_apps(idx=13, total=TOTAL)
print("  - Slide 13 (Data + Apps)")

# 14. Closing
slide_14_closing(idx=14, total=TOTAL)
print("  - Slide 14 (Closing)")

# =============== 保存 ===============
output = "c:/Users/user/Desktop/vibe/ppt1_caude/cv_history_deck.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Total slides: {len(prs.slides)}")
