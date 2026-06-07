"""结构验证: 解析 pptx 形状位置和文本,检测潜在溢出/重叠"""
from pptx import Presentation
from pptx.util import Emu

p = Presentation("c:/Users/user/Desktop/vibe/ppt1_caude/cv_history_deck.pptx")
sw, sh = p.slide_width, p.slide_height
print(f"Slide size: {sw/914400:.2f} x {sh/914400:.2f} inches")
print(f"Total slides: {len(p.slides)}")
print("=" * 70)

for i, slide in enumerate(p.slides, 1):
    print(f"\n--- Slide {i} ---")
    n_shapes = len(slide.shapes)
    print(f"  shapes: {n_shapes}")
    # 检查每个形状
    issues = []
    texts = []
    for shp in slide.shapes:
        if shp.has_text_frame:
            for para in shp.text_frame.paragraphs:
                for run in para.runs:
                    t = run.text.strip()
                    if t:
                        texts.append(t[:50])
    # 边界检查
    for shp in slide.shapes:
        if shp.left is None: continue
        right = shp.left + (shp.width or 0)
        bottom = shp.top + (shp.height or 0)
        if right > sw + 100:
            issues.append(f"  ⚠ 形状 {shp.shape_type} 右溢出: right={right/914400:.2f}\"")
        if bottom > sh + 100:
            issues.append(f"  ⚠ 形状 {shp.shape_type} 下溢出: bottom={bottom/914400:.2f}\"")
        if shp.left < -100:
            issues.append(f"  ⚠ 形状 {shp.shape_type} 左溢出: left={shp.left/914400:.2f}\"")
    # 输出
    if texts:
        print(f"  文本片段:")
        for t in texts[:6]:
            print(f"    · {t}")
        if len(texts) > 6:
            print(f"    ... +{len(texts)-6} more")
    for issue in issues:
        print(issue)
print("\n" + "=" * 70)
print("Done.")
